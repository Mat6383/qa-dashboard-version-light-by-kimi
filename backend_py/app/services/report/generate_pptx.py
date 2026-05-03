"""Generate PPTX closure report."""

from __future__ import annotations

import io
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt


def generate_pptx_report(data: dict[str, Any]) -> bytes | None:
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    tf = title.text_frame
    tf.text = "ISTQB Closure Report"
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True

    summary = data.get("summary", {})
    body = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(4))
    tf2 = body.text_frame
    tf2.text = f"Total: {summary.get('total', 0)}\nPassed: {summary.get('passed', 0)}\nFailed: {summary.get('failed', 0)}"

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()
